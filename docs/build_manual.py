"""Generate the TalkEx Business user manual as a PPTX.

Produces `TalkEx_Business_User_Manual.pptx` in the same directory. Uses
python-pptx for layout and PIL to draw the module-icon tiles so we don't
need external image assets.
"""

from __future__ import annotations

import io
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Design tokens — pulled from the frontend Tailwind palette so the manual
# looks like the product.
# ---------------------------------------------------------------------------

PRIMARY = RGBColor(0x25, 0x63, 0xEB)       # blue-600
PRIMARY_DARK = RGBColor(0x1E, 0x40, 0xAF)  # blue-800
PRIMARY_LIGHT = RGBColor(0xDB, 0xEA, 0xFE) # blue-100
INK = RGBColor(0x11, 0x18, 0x27)           # slate-900
INK_SOFT = RGBColor(0x47, 0x55, 0x69)      # slate-600
INK_MUTED = RGBColor(0x94, 0xA3, 0xB8)     # slate-400
BG = RGBColor(0xF8, 0xFA, 0xFC)            # slate-50
CARD = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)        # slate-200
GREEN = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUTPUT = Path(__file__).parent / "TalkEx_Business_User_Manual.pptx"

# ---------------------------------------------------------------------------
# PIL helpers — render UI mockup PNGs on the fly and stream them into the deck.
# ---------------------------------------------------------------------------

def _load_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    """Try common Windows fonts, fall back to PIL default."""
    for name in ("segoeuib.ttf" if bold else "segoeui.ttf",
                 "arialbd.ttf" if bold else "arial.ttf",
                 "DejaVuSans-Bold.ttf" if bold else "DejaVuSans.ttf"):
        try:
            return ImageFont.truetype(name, size)
        except OSError:
            continue
    return ImageFont.load_default()


def _rgb(color: RGBColor) -> tuple[int, int, int]:
    return (color[0], color[1], color[2])


def draw_app_shell(
    body_title: str,
    body_lines: list[tuple[str, str]],
    active_nav: str,
    right_widget: str | None = None,
) -> io.BytesIO:
    """Draw a shell mockup: dark sidebar + white main pane with title/lines.

    body_lines is a list of (label, value) pairs rendered as key-value rows.
    right_widget is optional short text placed as a badge in the header.
    """
    W, H = 1200, 720
    img = Image.new("RGB", (W, H), _rgb(BG))
    d = ImageDraw.Draw(img)

    # Sidebar
    sidebar_w = 220
    d.rectangle([0, 0, sidebar_w, H], fill=_rgb(INK))

    logo_font = _load_font(18, bold=True)
    d.text((22, 22), "TalkEx", font=logo_font, fill=(255, 255, 255))
    d.text((100, 22), "Business", font=logo_font, fill=(120, 200, 255))

    nav_groups = [
        ("Overview", ["Dashboard", "Channels"]),
        ("Messaging", ["Contacts", "Templates", "Campaigns", "Conversations"]),
        ("Automation & Dev", ["Automation", "Developers", "Webhooks", "Analytics", "Logs"]),
        ("Account", ["Billing", "Wallet", "Support", "Settings"]),
    ]

    label_font = _load_font(10, bold=True)
    item_font = _load_font(12)
    y = 70
    for group, items in nav_groups:
        d.text((22, y), group.upper(), font=label_font, fill=(148, 163, 184))
        y += 20
        for item in items:
            highlight = (item == active_nav)
            if highlight:
                d.rounded_rectangle([12, y - 4, sidebar_w - 12, y + 22], radius=6, fill=(37, 99, 235))
                d.text((22, y + 2), item, font=item_font, fill=(255, 255, 255))
            else:
                d.text((22, y + 2), item, font=item_font, fill=(203, 213, 225))
            y += 28
        y += 8

    # Header bar
    d.rectangle([sidebar_w, 0, W, 58], fill=(255, 255, 255))
    d.rectangle([sidebar_w, 57, W, 58], fill=_rgb(BORDER))
    d.text((sidebar_w + 24, 20), "One Platform. Multiple Messaging Channels.", font=_load_font(12), fill=_rgb(INK_SOFT))
    # bell + avatar
    d.ellipse([W - 130, 20, W - 110, 40], outline=_rgb(INK_SOFT), width=2)
    if right_widget:
        d.rounded_rectangle([W - 108, 22, W - 32, 40], radius=8, fill=(219, 234, 254))
        d.text((W - 102, 25), right_widget, font=_load_font(10, bold=True), fill=_rgb(PRIMARY_DARK))

    # Main pane
    d.rectangle([sidebar_w + 32, 96, W - 32, H - 32], fill=(255, 255, 255), outline=_rgb(BORDER), width=1)

    d.text((sidebar_w + 56, 116), body_title, font=_load_font(22, bold=True), fill=_rgb(INK))

    yy = 168
    key_font = _load_font(11, bold=True)
    val_font = _load_font(12)
    for k, v in body_lines:
        d.text((sidebar_w + 56, yy), k, font=key_font, fill=_rgb(INK_MUTED))
        d.text((sidebar_w + 200, yy), v, font=val_font, fill=_rgb(INK))
        yy += 30

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


# ---------------------------------------------------------------------------
# PPTX helpers
# ---------------------------------------------------------------------------

def _add_text_box(
    slide, left, top, width, height, text: str,
    size: int = 14, bold: bool = False, color: RGBColor = INK,
    align: PP_ALIGN = PP_ALIGN.LEFT, anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> None:
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    # Split on newlines so paragraphs stay separate.
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color


def _fill_bg(slide, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Push to back
    spTree = shape._element.getparent()
    spTree.remove(shape._element)
    spTree.insert(2, shape._element)


def _add_rect(slide, left, top, width, height, fill: RGBColor,
              line: RGBColor | None = None) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)


def _add_pill(slide, left, top, text: str, fill: RGBColor, text_color: RGBColor = INK) -> None:
    """Small rounded pill with centered text."""
    width = Inches(1.4)
    height = Inches(0.32)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.5
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = text_color


def new_slide(prs: Presentation, layout_idx: int = 6):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


# ---------------------------------------------------------------------------
# Slide templates
# ---------------------------------------------------------------------------

def slide_cover(prs: Presentation) -> None:
    s = new_slide(prs)
    _fill_bg(s, INK)

    # decorative gradient bar
    _add_rect(s, 0, 0, SLIDE_W, Inches(0.3), PRIMARY)

    _add_text_box(s, Inches(0.8), Inches(1.3), Inches(11.7), Inches(1),
                  "TalkEx Business", size=54, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    _add_text_box(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(0.6),
                  "One Platform. Multiple Messaging Channels.",
                  size=22, color=RGBColor(0x93, 0xC5, 0xFD))

    _add_text_box(s, Inches(0.8), Inches(3.6), Inches(11.7), Inches(0.5),
                  "User Manual", size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    _add_text_box(s, Inches(0.8), Inches(4.2), Inches(11.7), Inches(1.5),
                  "Step-by-step guide to every screen, sub-screen, and feature\n"
                  "of the TalkEx Business CPaaS dashboard.",
                  size=16, color=RGBColor(0xCB, 0xD5, 0xE1))

    _add_text_box(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.6),
                  "CoreAxis Ventures  •  v1.0  •  Generated by Claude",
                  size=12, color=RGBColor(0x94, 0xA3, 0xB8))


def slide_section_cover(prs: Presentation, number: str, title: str, subtitle: str) -> None:
    s = new_slide(prs)
    _fill_bg(s, BG)

    _add_rect(s, 0, Inches(2.5), Inches(0.4), Inches(2.5), PRIMARY)

    _add_text_box(s, Inches(1.2), Inches(2.5), Inches(2), Inches(1),
                  f"Section {number}", size=18, bold=True, color=PRIMARY)
    _add_text_box(s, Inches(1.2), Inches(3.2), Inches(11), Inches(1.5),
                  title, size=42, bold=True, color=INK)
    _add_text_box(s, Inches(1.2), Inches(4.4), Inches(11), Inches(2),
                  subtitle, size=16, color=INK_SOFT)


def slide_toc(prs: Presentation, entries: list[tuple[str, str]]) -> None:
    s = new_slide(prs)
    _fill_bg(s, CARD)

    _add_text_box(s, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
                  "What's in this manual", size=28, bold=True, color=INK)
    _add_text_box(s, Inches(0.6), Inches(1.05), Inches(12), Inches(0.4),
                  "The manual is organized so each module gets its own section — read straight through, or jump to what you need.",
                  size=12, color=INK_SOFT)

    # Two columns of TOC entries
    col_w = Inches(6)
    y_start = Inches(1.7)
    row_h = Inches(0.32)
    mid = (len(entries) + 1) // 2
    for i, (num, title) in enumerate(entries):
        col = 0 if i < mid else 1
        row = i if i < mid else i - mid
        left = Inches(0.6) + col * col_w
        top = y_start + row * row_h
        _add_text_box(s, left, top, Inches(0.6), row_h, num,
                      size=13, bold=True, color=PRIMARY)
        _add_text_box(s, left + Inches(0.7), top, Inches(5.2), row_h, title,
                      size=13, color=INK)


def slide_screen(
    prs: Presentation,
    module: str,
    screen: str,
    intro: str,
    steps: list[str],
    fields: list[tuple[str, str]] | None = None,
    tips: list[str] | None = None,
    mockup: io.BytesIO | None = None,
) -> None:
    """One screen documentation slide: title + how-to + optional mockup."""
    s = new_slide(prs)
    _fill_bg(s, CARD)

    # Header band
    _add_rect(s, 0, 0, SLIDE_W, Inches(1.1), INK)
    _add_text_box(s, Inches(0.6), Inches(0.22), Inches(12), Inches(0.4),
                  module.upper(), size=11, bold=True, color=RGBColor(0x93, 0xC5, 0xFD))
    _add_text_box(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.6),
                  screen, size=24, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    # Intro
    _add_text_box(s, Inches(0.6), Inches(1.35), Inches(12.1), Inches(0.7),
                  intro, size=13, color=INK_SOFT)

    # Layout: left = steps, right = mockup or fields+tips
    left_x = Inches(0.6)
    right_x = Inches(6.9)
    col_w = Inches(6.0)

    # Steps (left)
    y = Inches(2.15)
    _add_text_box(s, left_x, y, col_w, Inches(0.35),
                  "Steps", size=13, bold=True, color=PRIMARY)
    y += Inches(0.4)
    for i, step in enumerate(steps, 1):
        _add_text_box(s, left_x, y, Inches(0.4), Inches(0.4),
                      f"{i}.", size=12, bold=True, color=PRIMARY)
        _add_text_box(s, left_x + Inches(0.4), y, col_w - Inches(0.4), Inches(0.8),
                      step, size=12, color=INK)
        # Estimate row height by string length so long steps get more room
        lines = max(1, (len(step) // 55) + 1)
        y += Inches(0.28) * lines + Inches(0.06)

    # Right side
    if mockup is not None:
        s.shapes.add_picture(mockup, right_x, Inches(2.15), width=col_w, height=Inches(4.2))
    else:
        yy = Inches(2.15)
        if fields:
            _add_text_box(s, right_x, yy, col_w, Inches(0.35),
                          "Fields & actions", size=13, bold=True, color=PRIMARY)
            yy += Inches(0.4)
            for name, desc in fields:
                _add_text_box(s, right_x, yy, Inches(1.9), Inches(0.4),
                              name, size=11, bold=True, color=INK)
                _add_text_box(s, right_x + Inches(2), yy, col_w - Inches(2), Inches(0.6),
                              desc, size=11, color=INK_SOFT)
                lines = max(1, (len(desc) // 42) + 1)
                yy += Inches(0.24) * lines + Inches(0.06)
        if tips:
            yy += Inches(0.2)
            _add_text_box(s, right_x, yy, col_w, Inches(0.35),
                          "Tips", size=13, bold=True, color=AMBER)
            yy += Inches(0.4)
            for tip in tips:
                _add_text_box(s, right_x, yy, Inches(0.3), Inches(0.4),
                              "•", size=13, bold=True, color=AMBER)
                _add_text_box(s, right_x + Inches(0.3), yy, col_w - Inches(0.3), Inches(0.6),
                              tip, size=11, color=INK_SOFT)
                lines = max(1, (len(tip) // 50) + 1)
                yy += Inches(0.22) * lines + Inches(0.06)

    # Footer
    _add_text_box(s, Inches(0.6), Inches(7.1), Inches(12), Inches(0.3),
                  f"TalkEx Business — {module}", size=9, color=INK_MUTED)


# ---------------------------------------------------------------------------
# Actual manual content
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Cover
    slide_cover(prs)

    # TOC
    slide_toc(prs, [
        ("01", "Overview & Getting Started"),
        ("02", "Authentication (Sign in / Sign up)"),
        ("03", "Dashboard"),
        ("04", "Channels"),
        ("05", "Contacts"),
        ("06", "Templates"),
        ("07", "Campaigns"),
        ("08", "Conversations Inbox"),
        ("09", "Automation Rules"),
        ("10", "Developers — API Keys"),
        ("11", "Webhooks"),
        ("12", "Analytics"),
        ("13", "Activity Logs"),
        ("14", "Billing & Plans"),
        ("15", "Wallet"),
        ("16", "Support & Help"),
        ("17", "Settings"),
        ("18", "Notification Bell"),
        ("19", "Security & Best Practices"),
        ("20", "API Reference Cheat-Sheet"),
    ])

    # ------------------------- 01 Overview -------------------------
    slide_section_cover(prs, "01", "Overview & Getting Started",
        "Everything you need before you sign in for the first time — what TalkEx Business is, "
        "how to open the app, and how to keep your credentials safe.")

    slide_screen(prs, "Overview", "What is TalkEx Business?",
        "TalkEx Business is a CPaaS (Communication Platform as a Service) dashboard. It lets a business manage "
        "contacts, templates, campaigns, conversations, automation, billing, and analytics across multiple "
        "messaging channels (TalkEx first-party, WhatsApp Business, Telegram, Email, SMS, RCS) from one place.",
        steps=[
            "Every module is scoped to your account — you never see another business's data.",
            "Contacts, Templates, Campaigns, and Analytics are shared across every channel you enable.",
            "The Channel Connector is the only thing that differs per messaging platform.",
            "The 24-hour customer service window governs when free-form outbound replies are allowed.",
        ],
        fields=[
            ("Backend", "Go 1.26 · Gin · GORM · SQLite (dev) / Postgres (prod)"),
            ("Frontend", "React 19 · Vite · TypeScript · Tailwind CSS v4"),
            ("Auth", "JWT access + refresh rotation, bcrypt passwords"),
            ("API keys", "SHA-256 hashed, reveal-once plaintext"),
        ],
        tips=[
            "The whole product runs from a single compiled binary — no runtime install on deploy.",
            "Every browser request goes through the audit middleware, so you can debug from /logs.",
        ])

    slide_screen(prs, "Overview", "How to open the app",
        "TalkEx Business runs as two processes — a Go backend and a Vite frontend — that talk over HTTP.",
        steps=[
            "In one terminal: `go run ./cmd/server`   (backend starts on http://localhost:8080)",
            "In another: `cd frontend && npm install && npm run dev`   (frontend on http://localhost:5173)",
            "Open http://localhost:5173 in your browser.",
            "The Vite dev server proxies every request under /api to the Go backend automatically.",
            "First-time setup: copy .env.example to .env and set JWT_SECRET to a strong random string in production.",
        ],
        tips=[
            "Production refuses to start if JWT_SECRET is left at the default or is shorter than 32 chars.",
            "CORS_ORIGINS='*' is rejected in production; list your real frontend origin(s).",
        ])

    # ------------------------- 02 Auth -------------------------
    slide_section_cover(prs, "02", "Authentication",
        "How to register, sign in, sign out, and reset a forgotten password. Every dashboard action requires a valid session.")

    slide_screen(prs, "Auth", "Sign Up (Create an account)",
        "The register page lets a new business create its first admin account. That admin becomes the account owner.",
        steps=[
            "From the Sign in page, click 'Create one' at the bottom.",
            "Fill in Full Name (e.g. 'Priya Sharma').",
            "Choose your country code (+91 by default) and enter your mobile number.",
            "Enter your work email address.",
            "Choose a password of at least 8 characters. Click the eye icon to reveal.",
            "Tick 'I agree to Terms & Conditions & Privacy Policy'.",
            "Click 'Create account'. On success you are auto-logged-in and land on the Dashboard.",
        ],
        fields=[
            ("Send OTP (mobile)", "Sends a verification OTP; verify by entering the 4–6 digit code."),
            ("Send OTP (email)", "Same — an email OTP flow, useful when neither is yet trusted."),
            ("Social sign-up", "Google / Facebook / Apple / GitHub buttons wire to OAuth providers."),
        ],
        tips=[
            "Email is the login credential — it cannot be changed later; use a stable work address.",
            "OTP verification is optional at registration but recommended for account recovery.",
        ])

    slide_screen(prs, "Auth", "Sign In (Login)",
        "The login page accepts an email + password. On success, tokens are stored in localStorage and you are sent to the Dashboard.",
        steps=[
            "Go to http://localhost:5173/login (or click TalkEx logo → Sign in).",
            "Type your registered email.",
            "Type your password. Click the eye icon to reveal it while typing.",
            "(Optional) Tick 'Remember me for 30 days' for a longer session.",
            "Click 'Sign in'. On failure, a red error banner explains the reason.",
        ],
        fields=[
            ("Forgot password?", "Opens the reset flow — see the next slide."),
            ("Social login", "Google / Facebook / Apple / GitHub sign in when configured."),
            ("Create one", "Link at the bottom to the Sign up page."),
        ],
        tips=[
            "The API returns a generic 'invalid credentials' regardless of the reason (no email enumeration).",
            "Access tokens expire in 15 min; refresh tokens auto-rotate for 30 days.",
        ])

    slide_screen(prs, "Auth", "Forgot Password",
        "Requesting a reset email always returns 200 (so an attacker can't check whether an address is registered). "
        "If the email exists, a short-lived reset token is issued.",
        steps=[
            "On the Sign in page, click 'Forgot password?' next to the Password label.",
            "Type your registered email address.",
            "Click 'Send reset link'. You'll see a 'Check your email' confirmation.",
            "Open the reset email and click the link — it opens the reset form.",
            "Type your new password (min 8 chars), confirm, and click 'Reset password'.",
            "You're returned to Sign in — log in with the new password.",
        ],
        tips=[
            "The reset token expires in 30 minutes for safety.",
            "In dev the reset URL is logged to the server console, so no mail provider is needed to test.",
        ])

    slide_screen(prs, "Auth", "Sign Out",
        "You can sign out from either the header (top-right) or the sidebar footer.",
        steps=[
            "Header method: click the red door-arrow icon at the far right of the header.",
            "Sidebar method: click the sign-out icon next to 'by CoreAxis Ventures' in the sidebar footer.",
            "You are redirected to /login and both access + refresh tokens are cleared.",
        ],
        tips=[
            "Signing out only clears this browser — your account and data remain untouched.",
            "If a refresh token is stolen, sign in on that same browser to invalidate the old session (rotation).",
        ])

    # ------------------------- 03 Dashboard -------------------------
    slide_section_cover(prs, "03", "Dashboard",
        "Your at-a-glance home screen. Shows wallet balance, contacts total, templates, and message activity — every card links to its module.")

    slide_screen(prs, "Dashboard", "Dashboard home",
        "The first thing you see after logging in. Loads real numbers via /analytics/summary + /wallet + /contacts + /templates.",
        steps=[
            "Wallet Balance card — click to open the Wallet page.",
            "Total Contacts card — click to open the Contacts page.",
            "Message Templates card — click to open Templates.",
            "Messages Sent, Messages Received, Delivery Rate cards — pulled from /analytics/summary.",
            "The Message Activity area shows a placeholder note; use the Analytics page for the real chart.",
        ],
        fields=[
            ("Wallet Balance", "Formatted in INR by default (₹1,234.56)."),
            ("Delivery Rate", "Delivered + read / outbound × 100."),
            ("QualityBadge", "In the header — Green / Yellow / Red per WhatsApp quality tier."),
        ],
        tips=[
            "All cards auto-refresh on page load; navigate back to see fresh numbers.",
        ])

    # ------------------------- 04 Channels -------------------------
    slide_section_cover(prs, "04", "Channels",
        "Turn on the messaging channels you want to send from. Every enabled channel becomes available in Templates, Campaigns, and Conversations.")

    slide_screen(prs, "Channels", "Channels grid",
        "A responsive 3-column grid of channel tiles. Each tile shows the channel name, description, and an Enable/Disable button.",
        steps=[
            "Sidebar → 'Channels'.",
            "Locate the channel you want to enable (TalkEx, WhatsApp Business, Telegram, Email, SMS, RCS).",
            "Click 'Enable channel' — the button turns into 'Disable channel' and a green Enabled pill appears.",
            "For 'Coming Soon' channels, clicking enables a preview record so the row appears in your account, but sends aren't yet wired.",
            "Verified date shows when the channel was last enabled.",
        ],
        fields=[
            ("Enabled pill", "Green check — channel is active for your account."),
            ("Coming Soon pill", "Amber — the connector isn't implemented yet."),
            ("Disabled pill", "Grey — you have explicitly disabled the channel."),
        ],
        tips=[
            "Enabling a channel doesn't obligate you — you can toggle back off any time.",
            "Real send/receive integrations for TalkEx and WhatsApp arrive in later releases.",
        ])

    # ------------------------- 05 Contacts -------------------------
    slide_section_cover(prs, "05", "Contacts",
        "Your unified address book — shared across every channel. Add contacts one at a time, tag them, and see whose 24-hour window is open.")

    slide_screen(prs, "Contacts", "Contacts list",
        "The Contacts page shows every recipient in your account with phone, tags, consent status, and 24-hour window state.",
        steps=[
            "Sidebar → 'Contacts'.",
            "Use the search box (top-left) to filter by name, phone, email, or tag.",
            "Click 'Add Contact' (top-right) to open the create modal.",
            "Click the pencil icon on any row to edit.",
            "Click the trash icon → then 'Confirm' to delete.",
        ],
        fields=[
            ("Consent pill", "Green = opted in; Grey = no consent (can only be messaged inside the 24h window)."),
            ("Window pill", "Blue 'Open (24h)' = you can freely reply; empty = window closed, template required."),
            ("Tags chips", "Blue rounded chips — free-form labels for segmentation."),
        ],
        tips=[
            "Contacts are shared across ALL channels — one row per phone number.",
            "Bulk CSV import is on the roadmap; today create via API POST /contacts for scripting.",
        ])

    slide_screen(prs, "Contacts", "Add / Edit Contact modal",
        "The modal collects phone (required, immutable on edit), name, email, and comma-separated tags.",
        steps=[
            "Type the phone number in international format (e.g. +919876543210).",
            "Type the contact's name (optional but recommended).",
            "Type an email (optional).",
            "Type tags separated by commas (e.g. 'vip, mumbai, lead').",
            "Click 'Add Contact' — the modal closes and the new row prepends to the table.",
            "For edits, phone is disabled (grey) because it's the unique identifier.",
        ],
        tips=[
            "Tags are case-sensitive and free-form. Pick a small vocabulary for consistency.",
        ])

    # ------------------------- 06 Templates -------------------------
    slide_section_cover(prs, "06", "Templates",
        "Reusable message bodies with category tagging. Templates are required for outbound sends outside the 24-hour window.")

    slide_screen(prs, "Templates", "Templates list & create",
        "Templates are the pre-approved copy your campaigns and automation replies use. Each is tagged with a category that drives channel-level pricing rules.",
        steps=[
            "Sidebar → 'Templates'.",
            "Click 'New Template' (top-right).",
            "Enter Template Name (e.g. 'Welcome message').",
            "Pick a Category — Marketing / Utility / Authentication.",
            "Pick a Channel — one of your enabled channels.",
            "Type the Body. Use {{variable}} placeholders for merge fields (e.g. Hi {{name}}!).",
            "Click 'Save'. Status starts as 'Draft'; you (or an admin) can move it to 'Pending review' → 'Approved'.",
        ],
        fields=[
            ("Marketing", "Promotional / broadcast — highest per-message price on WhatsApp."),
            ("Utility", "Transactional (order confirmations, alerts) — lower price."),
            ("Authentication", "OTPs, verification codes — lowest price, strictest rules."),
        ],
        tips=[
            "Only Approved templates can be used outside the 24-hour service window.",
            "Body variables — anything in double-braces — are substituted per-recipient at send time.",
        ])

    # ------------------------- 07 Campaigns -------------------------
    slide_section_cover(prs, "07", "Campaigns",
        "Bulk-send flows that snapshot which template + which contacts to send + when to launch. Track sent/failed counters live.")

    slide_screen(prs, "Campaigns", "Campaigns list",
        "The list shows every campaign with status, template used, recipient count, and a tri-color progress bar (delivered / sent / failed).",
        steps=[
            "Sidebar → 'Campaigns'.",
            "Click 'New Campaign' to open the wizard.",
            "Row actions depend on status: Launch (draft), Cancel (running), Delete (terminal states).",
            "Delete requires an inline 'Confirm' click — no accidents.",
        ],
        fields=[
            ("Draft", "Not yet launched — you can still edit contacts/template."),
            ("Scheduled", "Will launch automatically at scheduled_at."),
            ("Running", "Fan-out in progress; cancellable."),
            ("Completed / Failed / Cancelled", "Terminal — Delete becomes available."),
        ])

    slide_screen(prs, "Campaigns", "Create Campaign wizard",
        "The wizard walks you through Name → Template → Recipients → Schedule.",
        steps=[
            "Enter a Campaign Name (e.g. 'Diwali promo').",
            "Pick a Template from the dropdown — only Approved templates on enabled channels appear.",
            "Multi-select recipients from your Contacts list. Use 'Select all' to include everyone.",
            "Toggle 'Schedule for later' if you don't want to send immediately.",
            "If scheduled, pick date + time.",
            "Click 'Create Campaign'. Status is 'Draft' (or 'Scheduled') until you Launch.",
        ],
        tips=[
            "Contact list is snapshotted at create-time; later contact edits don't change what a scheduled campaign sends.",
            "Zero-recipient campaigns are rejected with a validation error.",
        ])

    slide_screen(prs, "Campaigns", "Launch & monitor",
        "Launching starts a background goroutine that sends one message per recipient and updates counters live.",
        steps=[
            "From the Campaigns list, click 'Launch' on a Draft campaign row.",
            "Status flips to 'Running' immediately; the progress bar animates as sends complete.",
            "Refresh the page to see the latest sent_count / failed_count.",
            "On completion, a green in-app notification pops in the bell menu: 'Campaign completed: <name>. Sent X, failed Y of Z recipients.'",
            "The 'campaign.completed' webhook fires to every subscribed endpoint.",
            "To stop mid-run, click 'Cancel' — already-sent messages stand, further dispatches stop.",
        ])

    # ------------------------- 08 Conversations -------------------------
    slide_section_cover(prs, "08", "Conversations Inbox",
        "2-way chat inbox with automatic 24-hour window tracking. Left pane = inbox; right pane = thread + send box.")

    slide_screen(prs, "Conversations", "Inbox layout",
        "Every (contact, channel) pair is one conversation row. Newest-message-first order.",
        steps=[
            "Sidebar → 'Conversations'.",
            "The left pane lists every conversation with avatar, name, last-message time, unread count.",
            "Click a row to open its thread in the right pane. Unread count clears automatically.",
            "The 'Simulate inbound' button (dev helper) lets you pretend a contact replied — useful for testing the 24h window.",
        ],
        fields=[
            ("Unread badge", "Blue circle with count next to the contact name."),
            ("Window pill", "Header of the thread — Blue 'Open (24h)' or Grey 'Window closed'."),
            ("Contact meta", "Phone number + channel below the name."),
        ])

    slide_screen(prs, "Conversations", "Thread + Send box",
        "Right pane shows the message history and a send box that respects the 24-hour rule.",
        steps=[
            "Outbound (yours) messages appear right-aligned in primary blue with status underneath.",
            "Inbound messages appear left-aligned in white with a border.",
            "Type in the send box. It's disabled while the window is closed and shows the reason.",
            "Click 'Send' or press Cmd/Ctrl+Enter.",
            "The message appears in the thread immediately; the inbox row moves to the top.",
        ],
        tips=[
            "When the window is closed, you must send via a Template — the send box explains this and points to Templates.",
            "Every send fires the 'message.status' webhook so external systems can react.",
        ])

    # ------------------------- 09 Automation -------------------------
    slide_section_cover(prs, "09", "Automation Rules",
        "Fire an auto-reply when an inbound message matches your keywords. Fully case-insensitive; first-match-wins.")

    slide_screen(prs, "Automation", "Rules list",
        "Each rule shows name, keyword chips, reply preview, fired-count, and an Active/Paused toggle.",
        steps=[
            "Sidebar → 'Automation'.",
            "Click 'New Rule' to open the create modal.",
            "Toggle Active/Paused inline by clicking the green Active pill.",
            "Click the pencil to edit; the trash → Confirm to delete.",
        ],
        fields=[
            ("Fired", "How many inbound messages have matched this rule to date."),
            ("Match type", "Contains (default) / Starts with / Exact — under each keyword chip."),
        ])

    slide_screen(prs, "Automation", "Create / Edit Rule",
        "A rule has a name, one or more trigger keywords, a match type, and a reply.",
        steps=[
            "Enter a Rule Name (e.g. 'Pricing inquiry').",
            "Enter Trigger Keywords, comma-separated (e.g. 'price, cost, plans').",
            "Pick a Match Type: Contains / Starts with / Exact.",
            "Type the Auto-reply message (e.g. 'Our plans start at ₹0/mo. Visit /billing.').",
            "Optionally point the rule at a Template so the auto-reply uses approved copy.",
            "Tick 'Active' — the rule fires immediately on the next matching inbound.",
            "Click 'Create Rule'.",
        ],
        tips=[
            "Rules are checked in creation order — put your most-specific rules first.",
            "Because the reply is sent while the 24h window is open (a reply just arrived), no template is required.",
        ])

    # ------------------------- 10 Developers -------------------------
    slide_section_cover(prs, "10", "Developers — API Keys",
        "Server-to-server credentials for programmatic access to the same REST API the dashboard uses. Keys are SHA-256 hashed at rest; plaintext is shown exactly once.")

    slide_screen(prs, "Developers", "API keys list",
        "Every key shows a name, its 12-char display prefix, active/revoked status, and when it was last used.",
        steps=[
            "Sidebar → 'Developers'.",
            "Click 'Create API Key' (top-right) to open the create dialog.",
            "For an active key: click the ban icon to revoke.",
            "For a revoked key: click the trash icon to permanently delete.",
            "Delete is only offered after revoke — deleting an active key would leave holders with unexplained 401s.",
        ],
        fields=[
            ("Prefix", "The first 12 chars of the plaintext (e.g. 'txb_a1b2c3…')."),
            ("Last used", "Auto-updated on every API call that authenticates with the key."),
            ("Status", "Active (green) or Revoked (grey)."),
        ])

    slide_screen(prs, "Developers", "Create key → reveal secret",
        "When you create a key, the modal switches to a reveal panel that shows the plaintext exactly once.",
        steps=[
            "Click 'Create API Key'.",
            "Enter a Key Name (e.g. 'Production Server', 'CI', 'Staging').",
            "Click 'Create Key'.",
            "The dialog now shows an amber 'Copy this key now.' warning + the plaintext.",
            "Click the copy icon (right of the key) — checkmark confirms.",
            "Store the key in your secret manager (never in git). Click 'I've saved my key' to close.",
            "Use the key by sending it in the Authorization header: `Authorization: Bearer txb_...`",
        ],
        tips=[
            "Losing the plaintext means creating a new key — we cannot recover it.",
            "Every endpoint that accepts a JWT also accepts a valid API key (bearer form).",
        ])

    # ------------------------- 11 Webhooks -------------------------
    slide_section_cover(prs, "11", "Webhooks",
        "Outbound HTTP callbacks fired on platform events. Every delivery is signed with HMAC-SHA256 so your receiver can verify authenticity.")

    slide_screen(prs, "Webhooks", "Endpoints list",
        "Each endpoint row shows name, URL, subscribed events (as chips), last-fired time, and status.",
        steps=[
            "Sidebar → 'Webhooks'.",
            "Click 'New Endpoint' (top-right) to add a receiver.",
            "Click the chevron on any row to view its delivery log.",
            "Click the trash → 'Confirm' to remove.",
        ],
        fields=[
            ("Events subscribed", "Chips: inbound.message, message.status, campaign.completed, contact.created."),
            ("Last fired", "Timestamp of the most recent delivery attempt."),
            ("Status", "Active (green) or Paused (grey)."),
        ])

    slide_screen(prs, "Webhooks", "Create endpoint → verify signature",
        "Creating an endpoint returns a per-endpoint shared secret exactly once. Every POST to your URL includes X-TalkEx-Signature = HMAC-SHA256(secret, body).",
        steps=[
            "Click 'New Endpoint'.",
            "Enter a Name (e.g. 'CRM sync').",
            "Enter the Endpoint URL (must be public HTTPS; must return 2xx to ack).",
            "Tick the events you want (usually leave all four checked initially).",
            "Tick 'Active' and click 'Create Endpoint'.",
            "Copy the plaintext secret from the amber reveal panel — this is the only chance.",
            "In your receiver, verify: hmac_sha256(secret, request.body) == request.headers['X-TalkEx-Signature'].",
            "Use developers.ConstantTimeEqual to compare (avoids timing attacks).",
        ],
        tips=[
            "Deliveries have a 10-second HTTP timeout; slow receivers land in the failure log.",
            "Delivery log shows the full payload sent + the receiver's response code / error.",
        ])

    # ------------------------- 12 Analytics -------------------------
    slide_section_cover(prs, "12", "Analytics",
        "Aggregate KPIs, line chart of daily message volume, status donut, and channel-mix bars.")

    slide_screen(prs, "Analytics", "KPI cards + charts",
        "Everything is scoped to your account. Data comes from the messages/conversations/contacts/campaigns tables via JOINs.",
        steps=[
            "Sidebar → 'Analytics'.",
            "Top row: 4 KPI cards — Messages Sent, Messages Received, Delivery Rate, Open Windows.",
            "Second row: 4 more — Contacts, Active Campaigns, Total Messages, Active Channels.",
            "Line chart: dual-series Outbound/Inbound counts per day; date-range toggle 7d / 30d / 90d.",
            "Below: Outbound-by-Status donut (queued/sent/delivered/read/failed) with percentages.",
            "Bottom right: Conversations-by-Channel horizontal bars.",
        ],
        fields=[
            ("Delivery rate", "Delivered + read / outbound × 100."),
            ("Open windows", "Conversations with inbound in the last 24h."),
            ("Active campaigns", "Scheduled + running status."),
        ])

    # ------------------------- 13 Logs -------------------------
    slide_section_cover(prs, "13", "Activity Logs",
        "Every API request from your account — including failures — recorded with method, path, status, latency, IP, and error body.")

    slide_screen(prs, "Logs", "Filter, search, expand",
        "The Logs page is your first stop when something isn't working. Every request the backend saw is here.",
        steps=[
            "Sidebar (Automation & Dev group) → 'Logs'.",
            "Top: stat cards for Total Requests, Failed Requests, Success Rate.",
            "Search box filters by path (e.g. type '/contacts' to see only that endpoint).",
            "Method dropdown: All / GET / POST / PATCH / DELETE.",
            "'Failed only' checkbox: hide 2xx responses.",
            "Failed rows are clickable — click to expand the truncated error response body.",
            "Click Refresh (top-right) to reload.",
        ],
        tips=[
            "Method chips are color-coded — GET blue, POST green, PATCH amber, DELETE red.",
            "Latency > 500ms is a hint the underlying query needs an index.",
        ])

    # ------------------------- 14 Billing -------------------------
    slide_section_cover(prs, "14", "Billing & Plans",
        "Manage your subscription tier, see usage against your monthly quota, and view invoices.")

    slide_screen(prs, "Billing", "Current plan + tiers",
        "The Billing page shows a gradient 'Current plan' card at the top and three plan tiles below.",
        steps=[
            "Sidebar → 'Billing'.",
            "Current Plan card: shows tier name, monthly price, renewal date, and 'X / Y messages used' progress.",
            "Below: three plan tiles — Starter (Free), Growth (₹2,499 – Recommended), Scale (₹9,999).",
            "Click 'Switch to <plan>' on any tile that isn't your current plan.",
            "Confirmation is immediate: plan flips, an invoice is issued, and the message counter resets to 0.",
            "New invoice appears in the Invoices table at the bottom.",
        ],
        fields=[
            ("Included messages", "Free monthly quota — Starter 1k / Growth 15k / Scale 100k."),
            ("Overage", "Per-message charge after the quota is exhausted."),
            ("Recommended", "Growth is highlighted with a blue border + badge."),
        ])

    slide_screen(prs, "Billing", "Invoices",
        "Each plan change issues an invoice. Invoices are read-only from the UI today.",
        steps=[
            "Scroll to the 'Invoices' section at the bottom of the Billing page.",
            "Every row shows date, plan, billing period, amount (₹), status pill (paid/pending/failed).",
            "PDF download button is present but disabled until PDF generation is wired.",
        ],
        tips=[
            "Payment provider integration (Razorpay/Stripe) is planned; today plan switches are recorded but do not charge a card.",
        ])

    # ------------------------- 15 Wallet -------------------------
    slide_section_cover(prs, "15", "Wallet",
        "Your pre-paid message credit balance. Debits are logged as ledger rows with idempotency keys so retries never double-charge.")

    slide_screen(prs, "Wallet", "Balance + ledger",
        "The Wallet page shows the current balance card and the full transaction ledger newest-first.",
        steps=[
            "Sidebar → 'Wallet'.",
            "Big balance card at the top: current amount, currency (INR by default).",
            "Click 'Top up' to add credits (dev-only until a payment provider is wired).",
            "Ledger table below: type (credit/debit) chips, amount (green / red), balance_after, reference note.",
            "Idempotency-key column shows the client-supplied uniqueness token — safe to retry.",
        ],
        tips=[
            "Concurrent debits use SELECT FOR UPDATE inside a DB transaction — no race can drive the balance negative.",
            "A duplicate idempotency key returns the ORIGINAL transaction row; the balance is not changed.",
        ])

    # ------------------------- 16 Support -------------------------
    slide_section_cover(prs, "16", "Support & Help Center",
        "In-product help: FAQ accordion answering the most common questions, doc grid to jump to modules, and a ticket form for humans.")

    slide_screen(prs, "Support", "Help center + tickets",
        "Everything on one page — quick channels, FAQ, docs grid, and your tickets.",
        steps=[
            "Sidebar → 'Support'.",
            "Top: three quick-channel tiles — Email us, File a ticket, Docs & source.",
            "Middle: expandable FAQ accordion. Click any question to reveal the answer.",
            "Documentation grid: cards linking to Contacts, Templates, Campaigns, etc.",
            "Bottom: 'Your tickets' table — everything you filed via the 'Contact Support' button.",
        ])

    slide_screen(prs, "Support", "File a ticket",
        "Tickets accumulate for a human to answer; auto-reply / status updates come in a later release.",
        steps=[
            "Click 'Contact Support' (top-right).",
            "Enter a Subject (e.g. 'Webhook not firing for status updates').",
            "Pick a Priority: Low / Normal / High / Urgent.",
            "Describe the issue — include steps, expected vs actual, request IDs from Logs.",
            "Click 'Submit Ticket'. The row appears in your tickets table immediately.",
        ])

    # ------------------------- 17 Settings -------------------------
    slide_section_cover(prs, "17", "Settings",
        "Profile info, password change, business category. 2FA / team invites / danger zone arrive in a later release.")

    slide_screen(prs, "Settings", "Profile & password",
        "The Settings page is grouped into sections; only Profile + Change Password are wired today.",
        steps=[
            "Sidebar → 'Settings'.",
            "Profile section: edit Full Name and Business Category. Email is read-only.",
            "Click 'Save changes' — success banner appears.",
            "Change Password section: type Current, New (min 8), Confirm.",
            "Click 'Update password'. On success you stay logged in.",
        ],
        tips=[
            "Wrong current password returns a specific error inline (not a generic 'try again').",
            "Sessions do not invalidate on password change in this release — plan is to add it.",
        ])

    # ------------------------- 18 Notification bell -------------------------
    slide_section_cover(prs, "18", "Notification Bell",
        "Every actionable event lands in the in-app bell dropdown in the header. Badge count polls every 30 seconds.")

    slide_screen(prs, "Notification Bell", "Bell menu & mark read",
        "The bell icon lives at the top-right, next to your avatar. Red badge shows unread count.",
        steps=[
            "Click the bell icon in the header.",
            "Dropdown lists notifications newest first, unread items bolded with a dot.",
            "Type-coloured icons: Info (blue) / Success (green) / Warning (amber) / Error (red).",
            "Click any row — you're navigated to its `link` (e.g. /conversations) AND the row is marked read.",
            "Click 'Mark all read' at the top of the dropdown to clear the badge in one shot.",
            "Close by clicking anywhere outside the dropdown.",
        ],
        tips=[
            "Notifications are emitted server-side by other modules (inbound message, campaign completion, plan change).",
            "The badge polls every 30 seconds; hit refresh for immediate updates.",
        ])

    # ------------------------- 19 Security -------------------------
    slide_section_cover(prs, "19", "Security & Best Practices",
        "How TalkEx Business protects your data — and what you should do on your side.")

    slide_screen(prs, "Security", "How the platform protects you",
        "Multiple defensive layers.",
        steps=[
            "Passwords hashed with bcrypt (never stored in plaintext, never in logs).",
            "JWT access tokens (15-min lifetime) + refresh tokens (30-day, rotated on every use).",
            "Every DB query is owner-scoped — a query without `WHERE owner_id = ?` cannot leave the review pipeline.",
            "API keys and webhook secrets are stored SHA-256 hashed; plaintext is returned exactly once.",
            "Every request is captured in the audit log with method, path, status, latency, IP.",
            "Non-dev environments refuse to start with the default JWT_SECRET or CORS_ORIGINS='*'.",
            "Wallet debits run in a DB transaction with SELECT FOR UPDATE to prevent negative-balance races.",
        ])

    slide_screen(prs, "Security", "What you should do",
        "Best practices on the customer side.",
        steps=[
            "Rotate API keys quarterly — revoke old, create new, update your servers, delete old row.",
            "Verify every webhook: HMAC-SHA256(secret, body) MUST match X-TalkEx-Signature header. Use constant-time compare.",
            "Store secrets in a secret manager (Vault / AWS Secrets Manager / GCP Secret Manager) — NEVER in git.",
            "Turn on webhook TLS at your endpoint; TalkEx Business rejects http:// URLs.",
            "Prefer template-based sends outside the 24-hour window — free-form sends will 409.",
            "Review /logs weekly; any 5xx cluster is worth investigating in your app before it becomes a customer complaint.",
        ])

    # ------------------------- 20 API cheat-sheet -------------------------
    slide_section_cover(prs, "20", "API Reference Cheat-Sheet",
        "Every dashboard action is also available as a REST call. Bearer JWT or API key both accepted.")

    slide_screen(prs, "API", "Common endpoints",
        "Everything runs under http://localhost:8080 in dev; add your custom domain in prod.",
        steps=[
            "POST /auth/register  {email, password, full_name}  → user",
            "POST /auth/login     {email, password}             → {access_token, refresh_token}",
            "POST /auth/refresh   {refresh_token}               → new pair",
            "GET  /users/me                                     → current user",
            "POST /wallet/transactions  {type, amount, idempotency_key}  → txn (idempotent)",
            "POST /contacts       {phone_number, name, tags}    → contact",
            "POST /templates      {name, category, channel, body}  → template",
            "POST /campaigns      {name, template_id, contact_ids} → draft campaign",
            "POST /campaigns/:id/launch                         → running campaign + background sends",
            "POST /conversations/send  {contact_id, channel, body}  → outbound message",
        ],
        fields=[
            ("Auth header", "Authorization: Bearer <jwt_access_token OR txb_api_key>"),
            ("Content-Type", "application/json for POST/PATCH"),
            ("Errors", "422 validation, 401 auth, 403 forbidden, 404 not found, 409 conflict, 500 server"),
        ])

    slide_screen(prs, "API", "Developer/webhook endpoints",
        "For programmatic integrations.",
        steps=[
            "POST /api-keys        {name}  → {api_key, plaintext}  (plaintext once)",
            "POST /api-keys/:id/revoke     → revoked row",
            "POST /webhooks        {name, url, events, active}    → {endpoint, plaintext_secret}",
            "GET  /webhooks/:id/deliveries → last 50 delivery attempts",
            "PUT  /channels/:kind  {enabled, config}              → channel config",
            "POST /automation/rules {name, trigger_keywords, response_body} → rule",
            "GET  /analytics/summary                              → KPI blob",
            "GET  /analytics/timeseries?days=30                   → per-day series",
            "GET  /audit-logs?failed=true&method=POST&search=/contacts → filtered logs",
            "POST /notifications/read-all                         → clear badge",
        ])

    # Save
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Wrote {OUTPUT}  ({OUTPUT.stat().st_size:,} bytes, {len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
